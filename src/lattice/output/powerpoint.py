"""Export the argument outline to a teaching PowerPoint deck.

Each section of the outline becomes a slide. The thesis becomes the
title slide. Claim bullets become slide bullets. Source citations on
each claim are placed in the slide's speaker notes so the presenter
can cite live.

The generated deck is intentionally plain — black text, no
backgrounds, standard 16:9 — so the user can apply their own
institution's template via PowerPoint's "Reuse Slides / Apply Theme"
afterwards without fighting auto-generated styling.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from ..graph.models import (
    AuthorGraph,
    Claim,
    Cluster,
    SectionRole,
)


# Map our section roles to user-facing labels for the slide subtitle.
_ROLE_LABELS = {
    SectionRole.introduction: "Introduction",
    SectionRole.argumentative: "Argument",
    SectionRole.evidence_synthesis: "Evidence synthesis",
    SectionRole.methodological: "Methods",
    SectionRole.counterargument: "Counterargument",
    SectionRole.conclusion: "Conclusion",
    SectionRole.appendix: "Appendix",
}


def build_teaching_deck(
    graph: AuthorGraph,
    clusters: list[Cluster] | None = None,
    project_name: str | None = None,
) -> Presentation:
    """Build a python-pptx Presentation object from the graph.

    The deck has:
      1. Title slide — project name + thesis statement
      2. Outline slide — table of contents listing each section
      3. One body slide per section — title + bulleted claims
      4. Conclusion slide — summary bullets pulling from claims
         tagged ``mechanism`` or in conclusion-role sections.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    title = project_name or graph.project_name or "Untitled paper"
    thesis = (graph.thesis_statement or "").strip()

    # ── 1. Title slide ──
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = title
    if title_slide.placeholders[1].has_text_frame:
        sub = title_slide.placeholders[1].text_frame
        sub.text = thesis or "(no thesis statement)"
        for para in sub.paragraphs:
            for run in para.runs:
                run.font.size = Pt(20)

    # ── 2. Outline (table of contents) slide ──
    bullet_layout = prs.slide_layouts[1]
    toc_slide = prs.slides.add_slide(bullet_layout)
    toc_slide.shapes.title.text = "Outline"
    toc_body = _body_text_frame(toc_slide)
    if toc_body is not None:
        toc_body.clear()
        first = True
        for section in graph.sections:
            if section.role == SectionRole.references:
                continue
            label = section.title or section.section_id
            para = toc_body.paragraphs[0] if first else toc_body.add_paragraph()
            para.text = f"{section.section_id.replace('s.', '').upper()}.  {label}"
            para.level = 0
            for run in para.runs:
                run.font.size = Pt(22)
            first = False

    # ── 3. One body slide per section ──
    claims_by_id = {c.claim_id: c for c in graph.claims}
    for section in graph.sections:
        if section.role == SectionRole.references:
            continue
        if section.section_id == "s.thesis":
            continue
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = section.title or section.section_id

        body_tf = _body_text_frame(slide)
        if body_tf is None:
            continue
        body_tf.clear()

        # Pull claims attached to this section (skip the thesis claim
        # since it's on the title slide).
        section_claims = [
            claims_by_id[cid] for cid in section.claim_ids
            if cid in claims_by_id and cid != "cl.thesis"
        ]
        if not section_claims:
            para = body_tf.paragraphs[0]
            para.text = "(no claims)"
            for run in para.runs:
                run.font.size = Pt(20)
                run.font.italic = True
        else:
            first = True
            for claim in section_claims:
                para = body_tf.paragraphs[0] if first else body_tf.add_paragraph()
                para.text = (claim.statement or "").strip()
                para.level = 0
                for run in para.runs:
                    run.font.size = Pt(20)
                first = False

        # Speaker notes: list source citations + any mechanism /
        # scope_conditions notes from the claim payload, so the
        # presenter has supporting context off-slide.
        notes_lines: list[str] = []
        for claim in section_claims:
            sources = sorted({ev.source for ev in claim.evidence if ev.source})
            mech = (getattr(claim, "mechanism", None) or "").strip()
            line_parts: list[str] = [f"• {claim.statement}"]
            if sources:
                line_parts.append(f"  Sources: {', '.join(sources)}")
            if mech:
                line_parts.append(f"  Mechanism: {mech}")
            notes_lines.append("\n".join(line_parts))
        if notes_lines:
            slide.notes_slide.notes_text_frame.text = "\n\n".join(notes_lines)

        # Subtitle / role label below the title (small grey strip).
        role_label = _ROLE_LABELS.get(section.role, "")
        if role_label:
            _add_role_subtitle(slide, role_label)

    return prs


def _body_text_frame(slide):
    """Find the body placeholder's text_frame on a content layout.
    pptx slide layouts vary across themes, so we walk placeholders
    looking for one that isn't the title and has a text_frame."""
    title_id = slide.shapes.title.placeholder_format.idx if slide.shapes.title else None
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx == title_id:
                continue
            if ph.has_text_frame:
                return ph.text_frame
        except (AttributeError, ValueError):
            continue
    return None


def _add_role_subtitle(slide, role_label: str) -> None:
    """Drop a subtle role label under the slide title."""
    from pptx.util import Inches as _In
    from pptx.dml.color import RGBColor
    left = _In(0.5)
    top = _In(0.95)
    width = _In(8)
    height = _In(0.3)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.text = role_label
    p = tf.paragraphs[0]
    for run in p.runs:
        run.font.size = Pt(12)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)


def write_teaching_deck(
    graph: AuthorGraph,
    target_path: Path,
    clusters: list[Cluster] | None = None,
    project_name: str | None = None,
) -> Path:
    """Build the deck and write it to ``target_path``. Returns the
    final path (so callers don't have to recompute it)."""
    prs = build_teaching_deck(graph, clusters=clusters, project_name=project_name)
    target_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(target_path))
    return target_path
